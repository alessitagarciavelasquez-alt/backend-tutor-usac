import os  # Permite leer variables de entorno y manejar rutas del sistema de archivos
import io  # Permite trabajar con archivos en memoria sin necesidad de guardarlos en disco
import base64  # Convierte imágenes binarias a texto Base64 para enviarlas a la API de visión de IA
from flask import Flask, request, jsonify  # Flask crea el servidor web; request lee datos entrantes; jsonify convierte respuestas a JSON
from flask_cors import CORS  # Permite que el frontend en Netlify se comunique con este backend sin bloqueos de seguridad
from flask_sqlalchemy import SQLAlchemy  # ORM para manejar la base de datos SQLite sin escribir SQL directo
from flask_talisman import Talisman  # Aplica cabeceras de seguridad HTTP y fuerza el uso de HTTPS
from openai import OpenAI  # Cliente compatible con DeepSeek para enviar consultas a la inteligencia artificial

app = Flask(__name__)  # Crea la instancia principal del servidor Flask usando el nombre de este archivo

CORS(  # Configura los permisos de acceso cruzado para que Netlify pueda llamar a este backend
    app,  # Vincula la configuración CORS a la instancia de Flask
    resources={  # Define las rutas y permisos específicos que se aplicarán
        r"/*": {  # Aplica la configuración a absolutamente todas las rutas del servidor
            "origins": "*",  # Acepta peticiones desde cualquier dominio externo sin restricciones
            "methods": ["POST", "OPTIONS"],  # Solo permite los métodos POST y OPTIONS
            "allow_headers": ["Content-Type"]  # Solo acepta la cabecera Content-Type en las peticiones
        }
    }
)

Talisman(  # Inicializa el módulo de seguridad y lo vincula al servidor Flask
    app,  # Le pasa la instancia de Flask para aplicar las políticas de seguridad
    content_security_policy=None,  # Desactiva CSP para no bloquear scripts externos usados en el frontend
    force_https=True,  # Obliga a que todas las conexiones sean por HTTPS seguro
    strict_transport_security=False,  # Desactiva HSTS para evitar conflictos durante el despliegue en Render
    session_cookie_secure=False,  # Permite cookies sin requerir HTTPS en todas las circunstancias
    frame_options='ALLOWALL'  # Permite que la app pueda mostrarse dentro de iframes desde cualquier origen
)

basedir = os.path.abspath(os.path.dirname(__file__))  # Obtiene la ruta absoluta de la carpeta donde vive este archivo app.py
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///' + os.path.join(basedir, 'tutor_ai.db')  # Define la ubicación del archivo de base de datos SQLite en la misma carpeta del proyecto
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False  # Desactiva las notificaciones de cambios de SQLAlchemy para ahorrar memoria y recursos

db = SQLAlchemy(app)  # Crea el objeto de base de datos vinculado a la configuración de Flask

client = OpenAI(  # Crea el cliente que se conectará a la API de DeepSeek para procesar las consultas
    api_key=os.getenv("DEEPSEEK_API_KEY"),  # Lee la clave secreta de la API desde las variables de entorno de Render
    base_url="https://api.deepseek.com"  # Apunta el cliente al servidor de DeepSeek en lugar del de OpenAI por defecto
)

class Historial(db.Model):  # Define la estructura de la tabla 'historial' en la base de datos
    id = db.Column(db.Integer, primary_key=True)  # Identificador único autoincremental para cada registro
    tipo = db.Column(db.String(50))  # Guarda el tipo de solicitud: 'consulta', 'mapa', 'resumen', etc.
    respuesta = db.Column(db.Text)  # Guarda la respuesta completa que generó la IA

with app.app_context():  # Abre el contexto de Flask para ejecutar operaciones de base de datos al iniciar
    db.create_all()  # Crea las tablas en el archivo .db si todavía no existen en el disco

# ==============================================================
# FUNCIÓN CENTRAL: Extrae contenido de CUALQUIER tipo de archivo
# ==============================================================
def extraer_contenido(archivo, extension):  # Recibe el objeto del archivo y su extensión para decidir cómo procesarlo
    """Extrae texto o imagen Base64 de cualquier tipo de archivo soportado."""

    # ── TEXTO PLANO: TXT, MD, CSV, JSON, HTML, XML ─────────────
    if extension in ['.txt', '.md', '.csv', '.json', '.html', '.xml']:  # Verifica si es un archivo de texto simple que no necesita librería externa
        try:  # Intenta leer el archivo directamente como texto plano
            return archivo.read().decode('utf-8', errors='ignore'), None  # Decodifica los bytes a texto UTF-8 e ignora caracteres inválidos
        except Exception as e:  # Captura cualquier error de lectura sin detener el servidor
            return None, f"[Error leyendo archivo de texto: {str(e)}]"  # Retorna None y el mensaje de error

    # ── PDF ─────────────────────────────────────────────────────
    elif extension == '.pdf':  # Verifica si el archivo es un PDF
        try:  # Intenta extraer el texto del PDF usando pypdf sin necesidad de GPU
            from pypdf import PdfReader  # Importa el lector de PDF solo cuando se necesita para ahorrar memoria
            reader = PdfReader(io.BytesIO(archivo.read()))  # Lee el PDF desde memoria sin guardarlo en disco
            paginas = []  # Lista vacía para acumular el texto de cada página del documento
            for pagina in reader.pages:  # Itera sobre cada página del PDF
                texto = pagina.extract_text()  # Extrae el texto legible de la página actual
                if texto:  # Solo agrega la página si contiene texto extraíble
                    paginas.append(texto)  # Agrega el texto de la página a la lista acumuladora
            contenido = "\n".join(paginas)  # Une el texto de todas las páginas con saltos de línea
            if not contenido.strip():  # Si el PDF no tiene texto (puede ser un PDF escaneado como imagen)
                return None, "[El PDF no contiene texto extraíble. Puede ser un PDF escaneado. Intenta convertirlo a imagen PNG.]"  # Informa la limitación al usuario
            return contenido, None  # Retorna el texto completo extraído sin errores
        except Exception as e:  # Captura cualquier fallo durante la lectura del PDF
            return None, f"[Error al leer PDF: {str(e)}]"  # Retorna el mensaje de error detallado

    # ── WORD (.docx) ────────────────────────────────────────────
    elif extension == '.docx':  # Verifica si es un documento de Microsoft Word
        try:  # Intenta extraer el texto de todos los párrafos del documento Word
            from docx import Document  # Importa la librería python-docx solo cuando se necesita
            doc = Document(io.BytesIO(archivo.read()))  # Lee el documento Word desde memoria sin guardarlo en disco
            parrafos = [p.text for p in doc.paragraphs if p.text.strip()]  # Extrae solo los párrafos que tienen contenido real, ignorando los vacíos
            return "\n".join(parrafos), None  # Une todos los párrafos con saltos de línea y los retorna
        except Exception as e:  # Captura cualquier fallo al leer el documento Word
            return None, f"[Error al leer DOCX: {str(e)}]"  # Retorna el mensaje de error detallado

    # ── EXCEL (.xlsx / .xls) ────────────────────────────────────
    elif extension in ['.xlsx', '.xls']:  # Verifica si es una hoja de cálculo de Microsoft Excel
        try:  # Intenta leer todas las hojas y celdas del archivo Excel
            import openpyxl  # Importa la librería openpyxl solo cuando se necesita para leer Excel
            wb = openpyxl.load_workbook(io.BytesIO(archivo.read()), data_only=True)  # Carga el libro Excel desde memoria mostrando valores calculados en lugar de fórmulas
            todas_hojas = []  # Lista para acumular el contenido de todas las hojas del libro
            for nombre_hoja in wb.sheetnames:  # Itera sobre el nombre de cada hoja del libro Excel
                hoja = wb[nombre_hoja]  # Obtiene el objeto completo de la hoja actual
                filas = []  # Lista para acumular las filas de texto de esta hoja
                for fila in hoja.iter_rows(values_only=True):  # Itera sobre cada fila extrayendo solo los valores de las celdas
                    fila_texto = [str(c) for c in fila if c is not None]  # Convierte cada celda a texto ignorando las celdas vacías
                    if fila_texto:  # Solo agrega la fila si tiene al menos una celda con contenido
                        filas.append(" | ".join(fila_texto))  # Une las celdas con separador para simular una tabla de texto
                todas_hojas.append(f"[Hoja: {nombre_hoja}]\n" + "\n".join(filas))  # Agrega el nombre de la hoja seguido de sus filas al acumulador
            return "\n\n".join(todas_hojas), None  # Une todas las hojas con doble salto de línea y las retorna
        except Exception as e:  # Captura cualquier fallo al leer el archivo Excel
            return None, f"[Error al leer Excel: {str(e)}]"  # Retorna el mensaje de error detallado

    # ── POWERPOINT (.pptx) ──────────────────────────────────────
    elif extension == '.pptx':  # Verifica si es una presentación de Microsoft PowerPoint
        try:  # Intenta extraer el texto de todas las diapositivas y sus elementos
            from pptx import Presentation  # Importa la librería python-pptx solo cuando se necesita
            prs = Presentation(io.BytesIO(archivo.read()))  # Carga la presentación desde memoria sin guardarla en disco
            diapositivas = []  # Lista para acumular el texto de cada diapositiva
            for i, slide in enumerate(prs.slides, 1):  # Itera sobre cada diapositiva con su número de orden
                textos = []  # Lista para acumular el texto de los elementos de esta diapositiva
                for shape in slide.shapes:  # Itera sobre cada elemento de la diapositiva como títulos, cuadros de texto, etc.
                    if hasattr(shape, "text") and shape.text.strip():  # Verifica si el elemento tiene texto y no está vacío
                        textos.append(shape.text.strip())  # Agrega el texto del elemento a la lista
                if textos:  # Solo agrega la diapositiva si tiene al menos un elemento con texto
                    diapositivas.append(f"[Diapositiva {i}]\n" + "\n".join(textos))  # Agrega el número de diapositiva y su contenido
            return "\n\n".join(diapositivas), None  # Une todas las diapositivas con doble salto de línea y las retorna
        except Exception as e:  # Captura cualquier fallo al leer la presentación PowerPoint
            return None, f"[Error al leer PPTX: {str(e)}]"  # Retorna el mensaje de error detallado

    # ── IMÁGENES: JPG, PNG, WEBP, GIF ──────────────────────────
    elif extension in ['.jpg', '.jpeg', '.png', '.webp', '.gif']:  # Verifica si el archivo es una imagen de cualquiera de estos formatos
        try:  # Intenta convertir la imagen a Base64 para enviarla directamente al modelo de IA con capacidad de visión
            from PIL import Image  # Importa Pillow solo cuando se necesita para abrir y procesar la imagen
            img_bytes = archivo.read()  # Lee los bytes completos de la imagen desde el formulario web
            img = Image.open(io.BytesIO(img_bytes))  # Abre la imagen en memoria para verificar que es una imagen válida
            if img.mode != 'RGB':  # Si la imagen no está en modo RGB (puede ser RGBA con transparencia o paleta de colores)
                img = img.convert('RGB')  # Convierte la imagen a modo RGB para asegurar compatibilidad con la API de IA
            buffer = io.BytesIO()  # Crea un buffer en memoria para guardar la imagen procesada sin usar disco
            img.save(buffer, format='JPEG', quality=85)  # Guarda la imagen en el buffer en formato JPEG con calidad del 85%
            img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')  # Convierte los bytes de la imagen a texto Base64 compatible con la API
            return None, None, img_base64  # Retorna tres valores: texto(None), error(None), Base64 de la imagen
        except Exception as e:  # Captura cualquier fallo al procesar la imagen
            return None, f"[Error al procesar imagen: {str(e)}]", None  # Retorna el mensaje de error y None para el Base64

    # ── FORMATO NO SOPORTADO ────────────────────────────────────
    else:  # Si la extensión no coincide con ningún formato conocido por el sistema
        return None, f"[Formato '{extension}' no soportado. Formatos válidos: PDF, DOCX, XLSX, PPTX, TXT, CSV, JSON, MD, JPG, PNG, WEBP.]"  # Informa al usuario los formatos aceptados


# ==============================================================
# ENDPOINT PRINCIPAL: Recibe datos del frontend y responde con IA
# ==============================================================
@app.route('/procesar', methods=['POST'])  # Registra la ruta /procesar que solo acepta peticiones POST desde el frontend
def procesar():  # Función que se ejecuta cada vez que el frontend envía un formulario con datos

    try:  # Inicia el bloque de control de errores para capturar cualquier fallo durante el procesamiento

        tipo_solicitud = request.form.get('tipo', 'consulta')  # Lee el tipo de solicitud del formulario; usa 'consulta' si no viene ninguno
        texto_usuario = request.form.get('texto', '').strip()  # Lee el texto escrito o dictado por el usuario eliminando espacios innecesarios
        contenido_extraido = ""  # Variable que almacenará el texto extraído del archivo adjunto si viene uno
        imagen_base64 = None  # Variable que almacenará la imagen en Base64 si el usuario sube una imagen

        if 'file' in request.files and request.files['file'].filename != '':  # Verifica si el usuario adjuntó un archivo y que ese archivo tiene nombre válido
            archivo = request.files['file']  # Obtiene el objeto completo del archivo enviado desde el frontend
            extension = os.path.splitext(archivo.filename)[1].lower()  # Extrae la extensión del archivo y la convierte a minúsculas para comparación segura

            resultado = extraer_contenido(archivo, extension)  # Llama a la función central de extracción pasándole el archivo y su extensión

            if len(resultado) == 3:  # Si la función retornó 3 valores significa que el archivo es una imagen con Base64
                _, error, imagen_base64 = resultado  # Desempaqueta: texto ignorado, posible error, y el Base64 de la imagen
                if error:  # Si hubo un error al procesar la imagen
                    contenido_extraido = error  # Guarda el mensaje de error como contenido para informar al usuario en la respuesta
            else:  # Si retornó 2 valores es un archivo de texto normal como PDF, DOCX, TXT, etc.
                texto, error = resultado  # Desempaqueta el texto extraído y el posible mensaje de error
                contenido_extraido = texto if texto else (error or "")  # Usa el texto extraído o el mensaje de error si no se pudo leer

        system_msg = (  # Construye el mensaje de sistema que define el comportamiento y personalidad de Tutor AI
            "Tu nombre es exclusivamente 'Tutor AI'. Eres un experto de ingeniería de la USAC. "  # Define el nombre y rol académico del asistente de IA
            "Eres proactivo y autónomo. Si el usuario te realiza una pregunta simple o general, "  # Indica que debe responder de forma directa sin rodeos innecesarios
            "responde con texto Markdown limpio, sin inyectar fórmulas matemáticas previas. "  # Para preguntas simples no debe incluir fórmulas sin contexto
            "Si la consulta es técnica, proporciona tablas estructuradas en Markdown y ecuaciones "  # Para consultas técnicas exige tablas y ecuaciones bien formateadas
            "en notación LaTeX estándar. No menciones bajo ninguna circunstancia nombres antiguos de desarrollo."  # Prohíbe mencionar nombres de versiones anteriores del proyecto
        )

        # ── Construcción del prompt según el tipo de entrada recibida ────
        if imagen_base64:  # Si el usuario subió una imagen se activa el modo de visión del modelo de IA
            messages = [  # Construye la lista de mensajes con la imagen incrustada en formato Base64
                {"role": "system", "content": system_msg},  # Mensaje de sistema con las instrucciones de comportamiento de Tutor AI
                {
                    "role": "user",  # Mensaje del usuario que combina la imagen y su instrucción de texto
                    "content": [  # El contenido es una lista porque mezcla imagen y texto en un solo mensaje
                        {
                            "type": "image_url",  # Tipo especial de la API para enviar imágenes en formato Base64
                            "image_url": {  # Objeto que contiene la URL de datos con la imagen codificada
                                "url": f"data:image/jpeg;base64,{imagen_base64}"  # Formato estándar Data URL para imágenes Base64
                            }
                        },
                        {
                            "type": "text",  # Tipo texto para agregar la instrucción del usuario junto a la imagen
                            "text": texto_usuario if texto_usuario else "Describe y analiza detalladamente el contenido de esta imagen."  # Si no escribió nada solicita un análisis general de la imagen
                        }
                    ]
                }
            ]
        elif contenido_extraido:  # Si se extrajo texto de un documento adjunto construye el prompt combinado con contexto
            prompt_final = (  # Une el contenido del archivo con la instrucción del usuario en un solo prompt
                f"CONTEXTO DEL ARCHIVO SOPORTE:\n{contenido_extraido}\n\n"  # Inyecta el texto extraído del documento como contexto para la IA
                f"INSTRUCCIÓN DEL ESTUDIANTE: {texto_usuario}"  # Agrega la pregunta o instrucción escrita o dictada por el usuario
            )
            messages = [  # Lista de mensajes estándar sin imagen para texto y documentos
                {"role": "system", "content": system_msg},  # Mensaje de sistema con las instrucciones de Tutor AI
                {"role": "user", "content": prompt_final}  # Mensaje del usuario con contexto del archivo y su pregunta
            ]
        else:  # Si no hay archivo adjunto el prompt es únicamente el texto escrito o dictado por el usuario
            prompt_final = texto_usuario if texto_usuario else "Preséntate brevemente."  # Usa el texto del usuario o solicita una presentación breve si no escribió nada
            messages = [  # Lista de mensajes simple solo con texto del usuario
                {"role": "system", "content": system_msg},  # Mensaje de sistema con las instrucciones de comportamiento
                {"role": "user", "content": prompt_final}  # Mensaje del usuario con su consulta de texto puro
            ]

        completion = client.chat.completions.create(  # Envía la solicitud completa a la API de DeepSeek para obtener la respuesta de la IA
            model="deepseek-chat",  # Modelo de DeepSeek que procesará la consulta con soporte de texto y visión
            messages=messages,  # Lista de mensajes construida dinámicamente según el tipo de archivo recibido
            max_tokens=2048  # Limita la longitud de la respuesta a 2048 tokens para controlar tiempos y costos de API
        )

        respuesta = completion.choices[0].message.content  # Extrae el texto de la primera respuesta generada por el modelo de IA

        nuevo = Historial(tipo=tipo_solicitud, respuesta=respuesta)  # Crea el objeto de registro con el tipo de solicitud y la respuesta obtenida
        db.session.add(nuevo)  # Agrega el nuevo registro a la sesión activa de la base de datos pendiente de confirmación
        db.session.commit()  # Confirma y guarda definitivamente el registro en el archivo SQLite en disco

        return jsonify({"respuesta": respuesta})  # Devuelve la respuesta de la IA al frontend en formato JSON con la clave 'respuesta'

    except Exception as e:  # Captura cualquier error inesperado que ocurra durante todo el procesamiento anterior
        return jsonify({"error": str(e), "respuesta": f"Error interno: {str(e)}"}), 500  # Devuelve el detalle del error en JSON con código HTTP 500

if __name__ == '__main__':  # Solo se ejecuta si este archivo se corre directamente y no es importado por otro módulo
    port = int(os.environ.get("PORT", 10000))  # Lee el puerto asignado por Render desde las variables de entorno o usa 10000 como valor por defecto
    app.run(host='0.0.0.0', port=port)  # Inicia el servidor Flask escuchando en todas las interfaces de red en el puerto definido
if __name__ == '__main__':  # Verifica si este archivo se está ejecutando directamente y no siendo importado por otro módulo
    port = int(os.environ.get("PORT", 10000))  # Lee el puerto asignado por Render desde las variables de entorno o usa 10000 como valor por defecto
    app.run(host='0.0.0.0', port=port)  # Inicia el servidor Flask escuchando en todas las interfaces de red en el puerto definido
